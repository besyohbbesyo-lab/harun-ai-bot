# document_plugin.py
import io
from datetime import datetime
from pathlib import Path

import requests
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt


class DocumentCreator:
    def __init__(self):
        self.output_dir = Path.home() / "Desktop"

    def get_image(self, slide_index: int):
        """Picsum'dan yüksek kaliteli fotoğraf çek"""
        try:
            url = f"https://picsum.photos/seed/{slide_index}/800/450"
            headers = {"User-Agent": "Mozilla/5.0"}
            response = requests.get(url, headers=headers, timeout=15)
            if response.status_code == 200:
                return response.content
            return None
        except Exception:
            return None

    def add_image_to_slide(self, slide, slide_index: int):
        """Slayta görsel ekle - sağ tarafa"""
        try:
            image_bytes = self.get_image(slide_index)
            if image_bytes:
                image_stream = io.BytesIO(image_bytes)
                slide.shapes.add_picture(
                    image_stream, Inches(6), Inches(1.5), Inches(3.5), Inches(4.5)
                )
                return True
            return False
        except Exception as e:
            print(f"Gorsel eklenemedi: {e}")
            return False

    def create_presentation(self, title: str, slides_content: list) -> str:
        """PowerPoint sunum oluştur"""
        try:
            prs = Presentation()

            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = "Harun AI tarafindan hazirlandi"
            self.add_image_to_slide(slide, 1)

            for i, content in enumerate(slides_content):
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)

                baslik = content.get("baslik", f"Slayt {i+1}")
                icerik = content.get("icerik", "")

                slide.shapes.title.text = baslik
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.text = icerik

                self.add_image_to_slide(slide, i + 2)

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{title.replace(' ', '_')[:40]}_{timestamp}.pptx"
            filepath = self.output_dir / filename
            prs.save(filepath)

            return str(filepath)

        except Exception as e:
            return f"Hata: {e}"

    def create_word_document(self, title: str, content: str) -> str:
        """Word doküman oluştur"""
        try:
            doc = Document()

            doc.add_heading(title, 0)
            doc.add_paragraph(f"Tarih: {datetime.now().strftime('%d.%m.%Y %H:%M')}")
            doc.add_paragraph("Harun AI tarafindan hazirlandi")
            doc.add_paragraph("")

            paragraphs = content.split("\n")
            for para in paragraphs:
                if para.strip():
                    if para.startswith("#"):
                        doc.add_heading(para.replace("#", "").strip(), 1)
                    else:
                        doc.add_paragraph(para)

            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{title.replace(' ', '_')[:40]}_{timestamp}.docx"
            filepath = self.output_dir / filename
            doc.save(filepath)

            return str(filepath)

        except Exception as e:
            return f"Hata: {e}"
